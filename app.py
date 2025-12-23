"""
AI Lesson Plan Generator - ADAPTIVE
Works with ANY book, ANY language, ANY subject
"""

import streamlit as st
import base64
import re
import time
from groq import Groq
from typing import List, Dict, Tuple, Optional
import json
import numpy as np
import warnings

warnings.filterwarnings('ignore')

import logging

logging.getLogger('torch').setLevel(logging.ERROR)

# ============================================================
# PAGE CONFIG
# ============================================================

st.set_page_config(
    page_title="AI Lesson Generator v9.1",
    layout="wide",
    page_icon="🎓"
)


# ============================================================
# CACHED LOADERS
# ============================================================

@st.cache_resource(show_spinner="Loading embeddings...")
def load_embeddings():
    from langchain_huggingface import HuggingFaceEmbeddings
    return HuggingFaceEmbeddings(
        model_name="sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2",
        model_kwargs={'device': 'cpu'},
        encode_kwargs={'normalize_embeddings': True}
    )


@st.cache_resource(show_spinner=False)
def load_chromadb():
    import chromadb
    from chromadb.config import Settings
    return chromadb.Client(Settings(anonymized_telemetry=False, is_persistent=False))


@st.cache_resource(show_spinner=False)
def load_reranker():
    from flashrank import Ranker
    return Ranker(model_name="ms-marco-MiniLM-L-12-v2")


@st.cache_resource(show_spinner=False)
def load_rouge():
    from rouge_score import rouge_scorer
    return rouge_scorer.RougeScorer(['rouge1', 'rouge2', 'rougeL'], use_stemmer=True)


# ============================================================
# UNIVERSAL FILE EXTRACTION
# ============================================================

def extract_pdf(file) -> str:
    """Extract text from PDF"""
    import fitz

    pdf_bytes = file.getvalue()
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")

    text = ""
    for page_num, page in enumerate(doc):
        page_text = page.get_text("text")
        if page_text and len(page_text.strip()) > 20:
            text += f"\n\n--- Page {page_num + 1} ---\n\n"
            text += page_text

    doc.close()

    # Clean up
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r' {2,}', ' ', text)

    return text


def extract_docx(file) -> str:
    """Extract text from DOCX"""
    from docx import Document
    import io

    doc = Document(io.BytesIO(file.getvalue()))

    text = ""
    for para in doc.paragraphs:
        if para.text.strip():
            text += para.text + "\n\n"

    # Extract tables
    for table in doc.tables:
        text += "\n[TABLE]\n"
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            text += row_text + "\n"
        text += "[/TABLE]\n\n"

    return text


def extract_pptx(file) -> str:
    """Extract text from PowerPoint"""
    from pptx import Presentation
    import io

    prs = Presentation(io.BytesIO(file.getvalue()))

    text = ""
    for slide_num, slide in enumerate(prs.slides, 1):
        text += f"\n\n--- Slide {slide_num} ---\n\n"

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text + "\n"

    return text


def extract_txt(file) -> str:
    """Extract text from TXT file"""
    content = file.getvalue()

    # Try different encodings
    for encoding in ['utf-8', 'cp1251', 'latin-1']:
        try:
            return content.decode(encoding)
        except:
            continue

    return content.decode('utf-8', errors='ignore')


def extract_image(file, client: Groq) -> str:
    """Extract text from image using Groq Vision"""

    b64 = base64.b64encode(file.getvalue()).decode('utf-8')

    filename = file.name.lower()
    if filename.endswith('.png'):
        mime = "image/png"
    elif filename.endswith('.gif'):
        mime = "image/gif"
    elif filename.endswith('.webp'):
        mime = "image/webp"
    else:
        mime = "image/jpeg"

    try:
        response = client.chat.completions.create(
            messages=[{
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": """Extract ALL text from this image. 
- Keep all formulas exactly as they appear
- Preserve structure and formatting
- Include all numbers and symbols
- If there's a table, format it clearly"""
                    },
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:{mime};base64,{b64}"}
                    }
                ]
            }],
            model="llama-3.2-11b-vision-preview",
            temperature=0.1,
            max_tokens=4000
        )
        return response.choices[0].message.content
    except Exception as e:
        st.error(f"Vision API error: {e}")
        return ""


def extract_epub(file) -> str:
    """Extract text from EPUB"""
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
        import io

        book = epub.read_epub(io.BytesIO(file.getvalue()))

        text = ""
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), 'html.parser')
                text += soup.get_text() + "\n\n"

        return text
    except ImportError:
        st.warning("Install ebooklib:  pip install ebooklib beautifulsoup4")
        return ""


def extract_html(file) -> str:
    """Extract text from HTML"""
    try:
        from bs4 import BeautifulSoup

        content = file.getvalue().decode('utf-8', errors='ignore')
        soup = BeautifulSoup(content, 'html. parser')

        for script in soup(["script", "style"]):
            script.decompose()

        return soup.get_text()
    except ImportError:
        st.warning("Install beautifulsoup4: pip install beautifulsoup4")
        return ""


def extract_file(file, client: Groq = None) -> Tuple[str, str]:
    """Universal file extractor"""
    filename = file.name.lower()

    if filename.endswith('.pdf'):
        return extract_pdf(file), "PDF"
    elif filename.endswith('. docx'):
        return extract_docx(file), "DOCX"
    elif filename.endswith('. doc'):
        try:
            return extract_docx(file), "DOC"
        except:
            st.warning("Old . doc format - please convert to .docx")
            return "", "DOC"
    elif filename.endswith('.pptx'):
        return extract_pptx(file), "PPTX"
    elif filename.endswith('.txt'):
        return extract_txt(file), "TXT"
    elif filename.endswith('.epub'):
        return extract_epub(file), "EPUB"
    elif filename.endswith(('.html', '.htm')):
        return extract_html(file), "HTML"
    elif filename.endswith(('.png', '.jpg', '. jpeg', '.gif', '.webp')):
        if client:
            return extract_image(file, client), "IMAGE"
        else:
            st.error("Image extraction requires Groq client")
            return "", "IMAGE"
    else:
        st.warning(f"Unsupported format: {filename}")
        return "", "UNKNOWN"


# ============================================================
# DOCUMENT ANALYSIS
# ============================================================

def analyze_document(client: Groq, text: str) -> Dict:
    """Analyze document to understand content"""
    total_len = len(text)
    sample_start = text[: 1500]
    sample_middle = text[total_len // 2:total_len // 2 + 1000] if total_len > 3000 else ""
    sample_end = text[-1000:] if total_len > 2000 else ""

    sample = f"{sample_start}\n.. .\n{sample_middle}\n...\n{sample_end}"

    try:
        response = client.chat.completions.create(
            messages=[{
                "role": "user",
                "content": f"""Analyze this document.  Return ONLY valid JSON: 

{{
    "language": "english" or "uzbek" or "russian" or "other",
    "subject": "specific subject (e.g., IELTS Reading, Mathematics, Physics)",
    "topic": "main topic of this document",
    "level": "beginner" or "intermediate" or "advanced",
    "key_concepts": ["concept1", "concept2", "concept3"],
    "has_examples": true or false
}}

DOCUMENT: 
{sample[: 2500]}

JSON:"""
            }],
            model="llama-3.1-8b-instant",
            temperature=0.1,
            max_tokens=400
        )

        result_text = response.choices[0].message.content.strip()

        if "```" in result_text:
            match = re.search(r'```(?:json)?\s*(.*? )\s*```', result_text, re.DOTALL)
            result_text = match.group(1) if match else "{}"

        return json.loads(result_text)

    except Exception as e:
        return {
            "language": "english",
            "subject": "General",
            "topic": "Document content",
            "level": "intermediate",
            "key_concepts": [],
            "has_examples": False
        }


# ============================================================
# CHUNKING & VECTOR STORE
# ============================================================

def chunk_text(text: str, max_size: int = 600) -> List[str]:
    """Smart chunking"""
    paragraphs = re.split(r'\n\s*\n', text)

    chunks = []
    current = ""

    for para in paragraphs:
        para = para.strip()
        if not para or len(para) < 30:
            continue

        if len(current) + len(para) < max_size:
            current = f"{current}\n\n{para}" if current else para
        else:
            if current and len(current) > 50:
                chunks.append(current.strip())
            current = para

    if current and len(current) > 50:
        chunks.append(current.strip())

    return chunks if chunks else [text[: max_size]]


class VectorStore:
    def __init__(self):
        self.embeddings = load_embeddings()
        self.client = load_chromadb()
        self.collection = None
        self.docs = []

    def create(self, texts: List[str]) -> Tuple[int, float]:
        start = time.time()
        self.docs = texts

        try:
            self.client.delete_collection("docs")
        except:
            pass

        self.collection = self.client.create_collection(
            name="docs",
            metadata={"hnsw:space": "cosine"}
        )

        embeddings = self.embeddings.embed_documents(texts)

        self.collection.add(
            embeddings=embeddings,
            documents=texts,
            ids=[f"doc_{i}" for i in range(len(texts))]
        )

        return len(texts), time.time() - start

    def search(self, query: str, k: int = 10) -> List[str]:
        if not self.collection:
            return []

        emb = self.embeddings.embed_query(query)
        results = self.collection.query(
            query_embeddings=[emb],
            n_results=min(k, len(self.docs))
        )

        return results['documents'][0] if results['documents'] else []


# ============================================================
# CONTEXT BUILDER
# ============================================================

def build_context(text: str, chunks: List[str], store: VectorStore,
                  analysis: Dict) -> Tuple[str, float]:
    """Build optimal context from document"""
    start = time.time()

    intro = text[:2500]

    topic = analysis.get('topic', '')
    concepts = analysis.get('key_concepts', [])

    search_queries = [
        topic,
        ' '.join(concepts[: 3]) if concepts else topic,
        f"{analysis.get('subject', '')} main concepts",
    ]

    retrieved = set()
    for q in search_queries:
        if q:
            results = store.search(q, k=8)
            retrieved.update(results)

    # Find examples
    examples = []
    example_patterns = [
        r'example', r'exercise', r'question', r'task', r'practice',
        r'misol', r'mashq', r'savol', r'topshiriq',
        r'пример', r'упражнение', r'задача', r'вопрос',
    ]

    for chunk in chunks:
        chunk_lower = chunk.lower()
        if any(p in chunk_lower for p in example_patterns):
            examples.append(chunk)

    # Rerank
    retrieved_list = list(retrieved)
    if retrieved_list:
        try:
            from flashrank import RerankRequest
            ranker = load_reranker()
            passages = [{"id": i, "text": d} for i, d in enumerate(retrieved_list)]
            request = RerankRequest(query=topic, passages=passages)
            reranked = ranker.rerank(request)
            retrieved_list = [retrieved_list[r["id"]] for r in reranked[: 12]]
        except:
            pass

    # Combine
    context_parts = [
        "=== DOCUMENT START ===",
        intro,
        "\n=== KEY CONTENT ===",
        "\n---\n".join(retrieved_list[: 10]),
    ]

    if examples:
        context_parts.extend([
            "\n=== EXAMPLES FROM DOCUMENT ===",
            "\n---\n".join(examples[: 5])
        ])

    context = "\n".join(context_parts)

    if len(context) > 20000:
        context = context[:20000]

    return context, time.time() - start


# ============================================================
# DYNAMIC PROMPT GENERATION
# ============================================================

def create_prompt(analysis: Dict, context: str) -> Tuple[str, str]:
    """Create dynamic prompt based on analysis"""
    lang = analysis.get('language', 'english').lower()
    subject = analysis.get('subject', 'General')
    topic = analysis.get('topic', 'Unknown')
    level = analysis.get('level', 'intermediate')
    concepts = analysis.get('key_concepts', [])

    concepts_str = ', '.join(concepts) if concepts else 'as found in source'

    if lang in ['uzbek', 'uz']:
        system = f"""Siz tajribali {subject} o'qituvchisisiz. 
45 daqiqalik dars rejasi tuzing. 

MUHIM QOIDALAR:
1.  FAQAT quyidagi manba materialidan foydalaning
2. Hech narsa o'ylab topmang - faqat manbadagi ma'lumotlar
3. Misollarni manbadan oling
4. Ta'riflarni manbadan ko'chiring
5. Agar manbada yo'q bo'lsa, yozmang"""

        user = f"""Mavzu: {topic}
Daraja: {level}
Asosiy tushunchalar: {concepts_str}

Quyidagi strukturada dars rejasi tuzing:
1.  MAQSAD (2 daq) - 4 ta aniq maqsad
2.  KIRISH (5 daq) - Mavzuga kirish
3. ASOSIY QISM (25 daq) - Tushunchalar, formulalar, qoidalar MANBADAN
4. AMALIYOT (10 daq) - Misollar va mashqlar MANBADAN
5. MUSTAQIL ISH (3 daq)
6. UYGA VAZIFA (3 daq)
7. XULOSA (2 daq)

MANBA: 
{context}"""

    elif lang in ['russian', 'ru']:
        system = f"""Вы опытный преподаватель {subject}. 
Создайте план урока на 45 минут. 

ВАЖНЫЕ ПРАВИЛА:
1. Используйте ТОЛЬКО материал из источника ниже
2. НЕ выдумывайте ничего - только информация из источника
3. Берите примеры из источника
4. Цитируйте определения из источника
5. Если чего-то нет в источнике, не включайте"""

        user = f"""Тема: {topic}
Уровень: {level}
Ключевые понятия: {concepts_str}

Структура урока: 
1. ЦЕЛИ (2 мин) - 4 конкретные цели
2. ВВЕДЕНИЕ (5 мин)
3. ОСНОВНАЯ ЧАСТЬ (25 мин) - Понятия, формулы, правила ИЗ ИСТОЧНИКА
4. ПРАКТИКА (10 мин) - Примеры и упражнения ИЗ ИСТОЧНИКА
5. САМОСТОЯТЕЛЬНАЯ РАБОТА (3 мин)
6. ДОМАШНЕЕ ЗАДАНИЕ (3 мин)
7. ИТОГИ (2 мин)

ИСТОЧНИК:
{context}"""

    else:  # English and others
        system = f"""You are an experienced {subject} teacher. 
Create a 45-minute lesson plan.

CRITICAL RULES:
1. Use ONLY the source material provided below
2. Do NOT invent or hallucinate any content
3. Take examples DIRECTLY from the source
4. Quote definitions from the source
5. If something is not in the source, do NOT include it
6. Use quotation marks "" for direct quotes"""

        user = f"""Topic: {topic}
Level: {level}
Key concepts: {concepts_str}

Create a lesson plan with this structure:
1. OBJECTIVES (2 min) - 4 specific learning objectives
2. INTRODUCTION (5 min) - Hook and context FROM SOURCE
3. MAIN CONTENT (25 min) - Concepts, rules, explanations FROM SOURCE
4. PRACTICE (10 min) - Examples and exercises FROM SOURCE
5. INDEPENDENT WORK (3 min)
6. HOMEWORK (3 min)
7. SUMMARY (2 min)

SOURCE MATERIAL:
{context}"""

    return system, user


# ============================================================
# GENERATION
# ============================================================

def generate_lesson(client: Groq, system: str, user: str,
                    source: str, num: int = 2) -> Tuple[str, Dict, float]:
    """Generate with self-consistency"""
    start = time.time()

    responses = []
    temps = [0.1, 0.2, 0.25][: num]

    for i, temp in enumerate(temps):
        try:
            resp = client.chat.completions.create(
                messages=[
                    {"role": "system", "content": system},
                    {"role": "user", "content": user}
                ],
                model="llama-3.3-70b-versatile",
                temperature=temp,
                max_tokens=8000
            )
            responses.append(resp.choices[0].message.content)
        except Exception as e:
            st.warning(f"Generation {i + 1}:  {str(e)[:50]}")
            time.sleep(2)

    if not responses:
        return "", {"error": "All failed"}, time.time() - start

    if len(responses) == 1:
        return responses[0], {"method": "single"}, time.time() - start

    # Verify
    scores = []
    for resp in responses:
        try:
            verify = client.chat.completions.create(
                messages=[{
                    "role": "user",
                    "content": f"""Rate how well this lesson plan uses the source material (0-100).
High score = content is from source.
Low score = content seems invented. 

LESSON (first 1500 chars):
{resp[:1500]}

SOURCE (first 1500 chars):
{source[:1500]}

Return ONLY a number 0-100:"""
                }],
                model="llama-3.1-8b-instant",
                temperature=0.1,
                max_tokens=10
            )

            text = verify.choices[0].message.content.strip()
            nums = re.findall(r'\b(\d{1,3})\b', text)
            score = float(nums[0]) if nums and int(nums[0]) <= 100 else 50.0
            scores.append(score)
        except:
            scores.append(50.0)

        time.sleep(0.5)

    best = int(np.argmax(scores))

    return responses[best], {
        "method": "self-consistency",
        "candidates": len(responses),
        "scores": scores,
        "selected": best
    }, time.time() - start


# ============================================================
# METRICS
# ============================================================

def calculate_metrics(source: str, generated: str, embeddings) -> Dict:
    """Calculate quality metrics"""
    from sklearn.metrics.pairwise import cosine_similarity

    # ROUGE
    rouge = load_rouge()
    r = rouge.score(source, generated)
    rouge_scores = {
        'rouge1': r['rouge1'].fmeasure * 100,
        'rouge2': r['rouge2'].fmeasure * 100,
        'rougeL': r['rougeL'].fmeasure * 100
    }

    # Semantic
    src_emb = embeddings.embed_query(source[: 3000])
    gen_emb = embeddings.embed_query(generated[:3000])
    semantic = float(cosine_similarity([src_emb], [gen_emb])[0][0]) * 100

    # Faithfulness
    gen_sents = [s.strip() for s in re.split(r'[.!?]\s+', generated) if len(s.strip()) > 15][: 25]
    src_sents = [s.strip() for s in re.split(r'[.! ?]\s+', source) if len(s.strip()) > 15][:50]

    faithfulness = 50.
    0
    grounded = 50.0

    if gen_sents and src_sents:
        gen_embs = embeddings.embed_documents(gen_sents)
        src_embs = embeddings.embed_documents(src_sents)

        faith_scores = []
        grounded_count = 0

        for ge in gen_embs:
            sims = cosine_similarity([ge], src_embs)[0]
            max_sim = float(max(sims))
            faith_scores.append(max_sim)
            if max_sim > 0.5:
                grounded_count += 1

        faithfulness = float(np.mean(faith_scores)) * 100
        grounded = grounded_count / len(gen_sents) * 100

    # Structure
    structure_markers = [
        'objective', 'introduction', 'main', 'practice', 'homework', 'summary',
        'maqsad', 'kirish', 'asosiy', 'amaliyot', 'uyga', 'xulosa',
        'цель', 'введение', 'основн', 'практика', 'домашн', 'итог'
    ]
    gen_lower = generated.lower()
    structure = sum(1 for m in structure_markers if m in gen_lower) / 6 * 100
    structure = min(structure, 100)

    # Examples
    examples = len(re.findall(r'(example|exercise|misol|mashq|пример|упражн)\s*\d*',
                              generated, re.IGNORECASE))

    # Overall
    overall = (
            rouge_scores['rouge1'] * 0.10 +
            rouge_scores['rouge2'] * 0.05 +
            semantic * 0.25 +
            faithfulness * 0.35 +
            structure * 0.15 +
            min(examples * 5, 10) * 1.0
    )

    if overall >= 75:
        grade = "A"
    elif overall >= 65:
        grade = "B"
    elif overall >= 55:
        grade = "C"
    elif overall >= 45:
        grade = "D"
    else:
        grade = "F"

    return {
        'overall': overall,
        'grade': grade,
        'rouge': rouge_scores,
        'semantic': semantic,
        'faithfulness': faithfulness,
        'grounded': grounded,
        'hallucination': 100 - grounded,
        'structure': structure,
        'examples': examples,
    }


# ============================================================
# PDF EXPORT
# ============================================================

def create_pdf(text: str) -> Optional[bytes]:
    """Create PDF from text"""
    from fpdf import FPDF

    if not text:
        return None

    clean = text
    for c, r in {"∅": "(empty)", "Ω": "Omega", "≤": "<=", "≥": ">=", "•": "-"}.items():
        clean = clean.replace(c, r)
    clean = re.sub(r'[#*]+', '', clean)
    clean = ''.join(c if ord(c) < 128 else '?' for c in clean)

    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_font('Helvetica', '', 10)

    for line in clean.split('\n'):
        try:
            pdf.multi_cell(0, 6, line.strip() if line.strip() else ' ')
        except:
            pass

    output = pdf.output()
    return bytes(output) if isinstance(output, bytearray) else output


# ============================================================
# MAIN APP
# ============================================================

st.title("🎓 AI Lesson Plan Generator v9.1")
st.caption("📚 Works with ANY book, ANY language, ANY subject")

# Sidebar
with st.sidebar:
    st.header("⚙️ Settings")
    api_key = st.text_input("Groq API Key", type="password")

    with st.expander("🔧 Advanced"):
        num_generations = st.slider("Generations", 1, 3, 2)
        show_analysis = st.checkbox("Show document analysis", value=True)
        show_timing = st.checkbox("Show timing", value=True)

    st.divider()
    st.markdown("### ✨ v9.1 Features")
    st.write("📄 PDF, DOCX, PPTX, TXT")
    st.write("🖼️ Images (OCR via AI)")
    st.write("🌍 Any language")
    st.write("📖 Any subject")
    st.write("🎯 Adaptive prompts")

# Check API key
if not api_key:
    st.info("👈 Enter your Groq API key to start")

    with st.expander("ℹ️ How to get API key"):
        st.markdown("""
        1. Go to [console.groq.com](https://console.groq.com)
        2. Create free account
        3. Generate API key
        4.  Paste it in the sidebar
        """)
    st.stop()

try:
    client = Groq(api_key=api_key)
except Exception as e:
    st.error(f"Error:  {e}")
    st.stop()

# Upload
uploaded = st.file_uploader(
    "📎 Upload files (PDF, DOCX, PPTX, Images, TXT)",
    type=["pdf", "docx", "doc", "pptx", "txt", "epub", "html", "htm",
          "png", "jpg", "jpeg", "gif", "webp"],
    accept_multiple_files=True
)

# Generate
if uploaded and st.button("🚀 Generate Lesson Plan", type="primary"):

    total_start = time.time()
    timings = {}

    # Step 1: Extract files
    all_text = ""
    for file in uploaded:
        with st.spinner(f"Extracting {file.name}..."):
            text, file_type = extract_file(file, client)

            if text:
                all_text += f"\n\n=== {file.name} ({file_type}) ===\n\n"
                all_text += text
                st.success(f"✅ {file.name}:  {len(text)} chars extracted")
            else:
                st.error(f"❌ {file.name}: extraction failed")

    if not all_text or len(all_text) < 100:
        st.error("Could not extract enough text from files")
        st.stop()

    progress = st.progress(0)
    status = st.empty()

    # Step 2: Analyze document
    status.text("🔍 Analyzing document...")
    analysis_start = time.time()
    analysis = analyze_document(client, all_text)
    timings['analysis'] = time.time() - analysis_start

    if show_analysis:
        with st.expander("📊 Document Analysis", expanded=True):
            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("Language", analysis.get('language', 'unknown').title())
                st.metric("Level", analysis.get('level', 'unknown').title())
            with col2:
                st.metric("Subject", analysis.get('subject', 'unknown'))
            with col3:
                st.metric("Topic", analysis.get('topic', 'unknown')[:30])

            if analysis.get('key_concepts'):
                st.write("**Key concepts:**", ", ".join(analysis['key_concepts']))

    progress.progress(0.20)

    # Step 3: Chunk & Index
    status.text("🧠 Building knowledge base...")
    chunks = chunk_text(all_text)
    store = VectorStore()
    n_chunks, t = store.create(chunks)
    timings['index'] = t
    st.success(f"✅ {n_chunks} chunks indexed ({t:.1f}s)")
    progress.progress(0.40)

    # Step 4: Build context
    status.text("🔎 Finding relevant content...")
    context, t = build_context(all_text, chunks, store, analysis)
    timings['context'] = t
    st.info(f"📚 Context built ({t:.1f}s, {len(context)} chars)")
    progress.progress(0.55)

    # Step 5: Generate
    status.text(f"✍️ Generating lesson plan ({num_generations} candidates)...")
    system_prompt, user_prompt = create_prompt(analysis, context)
    lesson_plan, info, t = generate_lesson(client, system_prompt, user_prompt,
                                           all_text, num=num_generations)
    timings['generation'] = t

    if not lesson_plan:
        st.error("❌ Generation failed")
        st.stop()

    st.success(f"✅ Generated ({t:.1f}s)")
    if info.get('scores'):
        st.caption(f"Scores: {info['scores']}, Best: #{info.get('selected', 0) + 1}")

    progress.progress(0.80)

    # Step 6: Metrics
    status.text("📊 Calculating quality...")
    embeddings = load_embeddings()
    metrics_start = time.time()
    metrics = calculate_metrics(all_text, lesson_plan, embeddings)
    timings['metrics'] = time.time() - metrics_start

    progress.progress(1.0)
    total = time.time() - total_start
    status.empty()
    progress.empty()

    # Results
    st.success(f"✅ Complete in {total:.1f} seconds!")

    if show_timing:
        parts = [f"{k}: {v:.1f}s" for k, v in timings.items()]
        st.caption(f"⏱️ {' | '.join(parts)}")

    # Quality
    st.markdown("## 📊 Quality Assessment")
    c1, c2 = st.columns([3, 1])
    with c1:
        score = metrics['overall']
        if score >= 65:
            st.success(f"### Score: {score:.1f}%")
        elif score >= 50:
            st.warning(f"### Score: {score:.1f}%")
        else:
            st.error(f"### Score:  {score:.1f}%")
    with c2:
        st.markdown(f"### Grade: {metrics['grade']}")

    # Metrics
    st.markdown("### 📈 Detailed Metrics")
    c1, c2, c3, c4, c5 = st.columns(5)

    with c1:
        st.metric("ROUGE-1", f"{metrics['rouge']['rouge1']:.1f}%")
        st.metric("ROUGE-2", f"{metrics['rouge']['rouge2']:.1f}%")
    with c2:
        st.metric("Semantic", f"{metrics['semantic']:.1f}%")
    with c3:
        st.metric("Faithfulness", f"{metrics['faithfulness']:.1f}%")
        st.metric("Grounded", f"{metrics['grounded']:.1f}%")
    with c4:
        st.metric("Hallucination", f"{metrics['hallucination']:.1f}%")
    with c5:
        st.metric("Structure", f"{metrics['structure']:.1f}%")
        st.metric("Examples", metrics['examples'])

    st.divider()

    # Lesson Plan
    st.markdown("## 📝 Generated Lesson Plan")
    st.text_area("Content", lesson_plan, height=500)

    st.divider()

    # Export
    st.markdown("## 💾 Export")
    st.download_button(
        "📄 Download TXT",
        lesson_plan,
        "lesson_plan.txt",
        "text/plain"
    )
    # c1, c2, c3 = st.columns(3)
    #
    # with c1:
    #     pdf_bytes = create_pdf(lesson_plan)
    #     if pdf_bytes:
    #         st.download_button(
    #             "📥 Download PDF",
    #             pdf_bytes,
    #             "lesson_plan.pdf",
    #             "application/pdf",
    #             type="primary"
    #         )
    #
    # with c2:
    #     export_data = {
    #         'lesson_plan': lesson_plan,
    #         'analysis': analysis,
    #         'metrics': metrics,
    #         'timings': timings
    #     }
    #     st.download_button(
    #         "📦 Download JSON",
    #         json.dumps(export_data, ensure_ascii=False, indent=2, default=str),
    #         "lesson_plan.json",
    #         "application/json"
    #     )
    #
    # with c3:
    #

# Footer
st.divider()
st.markdown(
    "<center>🎓 AI Lesson Plan Generator - Adaptive</center>",
    unsafe_allow_html=True
)